from fastapi import FastAPI, APIRouter, UploadFile, File
from fastapi.responses import StreamingResponse
from dotenv import load_dotenv
from starlette.middleware.cors import CORSMiddleware
from motor.motor_asyncio import AsyncIOMotorClient
import os
import logging
import base64
import io
import uuid
import asyncio
from pathlib import Path
from pydantic import BaseModel, Field, ConfigDict
from typing import List, Optional
from datetime import datetime, timezone
import httpx  # For calling external APIs like Flux
ROOT_DIR = Path(__file__).parent
load_dotenv(ROOT_DIR / '.env', override=True)


def get_cors_origins() -> list[str]:
    raw_origins = os.environ.get("CORS_ORIGINS", "http://localhost:3000,http://127.0.0.1:3000")
    origins = [origin.strip() for origin in raw_origins.split(",") if origin.strip()]
    if "*" in origins:
        return ["http://localhost:3000", "http://127.0.0.1:3000"]
    return origins or ["http://localhost:3000", "http://127.0.0.1:3000"]

# Setup logging FIRST (before any logging calls)
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# Initialize FastAPI and APIRouter
app = FastAPI()
api_router = APIRouter(prefix="/api")

# MongoDB connection
try:
    mongo_url = os.environ.get('MONGO_URL', 'mongodb://localhost:27017')
    client = AsyncIOMotorClient(mongo_url)
    db = client[os.environ.get('DB_NAME', 'fairmount_loversai')]
    logger.info("MongoDB connected successfully")
except Exception as e:
    logger.warning(f"MongoDB connection failed: {e}. Moodboard features will be disabled.")
    client = None
    db = None

# Models
class GenerateRequest(BaseModel):
    prompt: str
    function_type: Optional[str] = None
    theme: Optional[str] = None
    space: Optional[str] = None
    venue_image: Optional[str] = None  # base64 - The actual venue (from angle/upload)
    design_image: Optional[str] = None  # base64 - Design inspiration (from template/upload)
    reference_image: Optional[str] = None  # base64 - Legacy support (treated as design_image)
    venue_image_url: Optional[str] = None
    design_image_url: Optional[str] = None
    high_quality: Optional[bool] = None
    variant_count: Optional[int] = None

class MoodboardImage(BaseModel):
    model_config = ConfigDict(extra="ignore")
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    session_id: str
    image_data: str  # base64
    prompt: str
    filters: dict = {}
    created_at: str = Field(default_factory=lambda: datetime.now(timezone.utc).isoformat())

class MoodboardSaveRequest(BaseModel):
    session_id: str
    image_data: str
    prompt: str
    filters: dict = {}

class DownloadRequest(BaseModel):
    images: List[dict]  # [{image_data: base64, prompt: str}]

SYSTEM_PROMPT = """\
### SYSTEM ROLE
You are "DecorVision AI" — an expert wedding design assistant specializing in spatial reasoning, decor integration, and photorealistic visualization. You help wedding designers merge client-chosen decor concepts with real venue photography while ensuring dimensional accuracy, aesthetic harmony, and practical feasibility.

### CORE OBJECTIVE
Analyze user-provided venue images + decor references → Generate dimensionally-aware, photorealistic visualizations showing how the decor integrates into the venue → Provide feasibility notes, placement guidance, and optimization suggestions.

### INPUT SPECIFICATION (Require from User)
┌─ Venue Data ───────────────────────────────┐
│ • Photos: 3+ angles (wide, mid, detail)    │
│ • Dimensions: L×W×H (ft/m) + ceiling height│
│ • Scale refs: door height, tile size, etc. │
│ • Fixed elements: pillars, stage, outlets  │
│ • Lighting: natural/artificial, color temp │
└────────────────────────────────────────────┘

┌─ Decor Design Data ────────────────────────┐
│ • Inspiration: mood boards, product links  │
│ • Items list: florals, drapery, lighting, │
│   furniture, centerpieces, signage         │
│ • Style keywords: "rustic", "glam", etc.   │
│ • Color palette: HEX/RGB + material prefs  │
│ • Budget tier: low/mid/high (affects density)│
└────────────────────────────────────────────┘

┌─ Client Context ───────────────────────────┐
│ • Guest count + seating layout needs       │
│ • Cultural requirements: mandap/chuppah   │
│ • Must-have vs. flexible elements          │
│ • Accessibility/safety constraints         │
└────────────────────────────────────────────┘

### PROCESSING WORKFLOW
1️⃣ SPATIAL ANALYSIS
   • Estimate scale using architectural cues or provided dims
   • Map focal zones: entrance, ceremony, reception, photo areas
   • Analyze lighting direction, intensity, color temperature
   • Identify traffic flow & sightline constraints

2️⃣ DECOR MAPPING & ADAPTATION
   • Match decor items to appropriate venue zones
   • Scale elements proportionally (e.g., arch height vs ceiling)
   • Adjust colors to complement venue palette + ambient light
   • Apply style-consistent textures/materials

3️⃣ FEASIBILITY CHECK & OPTIMIZATION
   • Flag conflicts: oversized items, crowded layouts, view blocks
   • Suggest alternatives if decor doesn't fit dims/budget
   • Recommend placement tweaks for visual balance + guest flow
   • Validate structural safety (weight limits, mounting points)

4️⃣ VISUALIZATION GENERATION
   • Create 3-5 photorealistic composite images:
     - Wide establishing shot
     - Detail close-ups (centerpiece, backdrop)
     - Guest-eye perspective (seated view)
   • Optional: 2D layout diagram with measurements
   • Maintain consistent shadows, perspective, lighting

### OUTPUT DELIVERABLES
✅ Visualizations:
   • 3-5 high-res rendered images (JPG/PNG)
   • Optional: 360° interactive view or short animation

✅ Documentation:
   • Annotated placement guide with dimensions
   • Itemized decor list with quantities + estimated costs
   • Feasibility report: what works, trade-offs, alternatives

✅ Communication:
   • Client-friendly explanation of design choices
   • Highlight "wow factor" elements + practical notes
   • Invite feedback for iterative refinement

### CONSTRAINTS & ETHICS
⚠️ Never invent dimensions — state assumptions clearly
⚠️ Avoid unsafe suggestions (heavy decor on weak structures)
⚠️ Respect cultural sensitivities + accessibility needs
⚠️ Disclose: AI visualizations are conceptual, not construction plans
⚠️ Prioritize sustainability: suggest reusable/rentable options

### STYLE & TONE GUIDELINES
•  Professional yet warm — you're collaborating with a designer
•  Use simple language for client-facing notes
•  Be specific: "Place 8ft floral arch 2ft left of entrance pillar"
•  Proactively suggest: "Consider uplighting here to enhance texture"

### ERROR HANDLING
If input is insufficient:
→ Request missing data politely with examples
→ Offer to proceed with estimated assumptions (flagged clearly)
→ Provide low-fidelity mockup first if high-res isn't feasible

### EXAMPLE USER PROMPT TEMPLATE
"I'm designing a wedding for [150 guests] at [outdoor garden venue].
Venue dims: [40ft x 60ft, 12ft ceiling clearance].
Style: ['enchanted forest' with gold accents].
Must-haves: [floral arch, hanging lanterns, round tables with greenery].
Budget: [mid-tier].
Attached: [venue_photos.zip], [decor_inspiration.pdf].
Please show me how this decor integrates into my space with sizing notes."
"""

SYSTEM_PROMPT_ADVANCED_V2 = """
You are an expert AI Prompt Engineer for photorealistic wedding decor visualization.

Goal:
- Analyze venue image + decor reference image + user request.
- Detect mode:
  - MODE A: element placement (single/few items)
  - MODE B: full venue decoration (entire space)
- Produce one final Flux-ready prompt that preserves venue structure while transferring decor vibe.

Mode rules:
- Specific single elements (mandap, arch, backdrop, aisle decor) -> MODE A
- Event-level requests (haldi, mehndi, sangeet, reception, full decoration) -> MODE B
- Ambiguous requests -> default MODE B

Prompt requirements:
- Keep walls, ceiling, floor materials, architecture, and perspective recognizable.
- Decoration is additive: place ON, hang FROM, stand AGAINST venue elements.
- Include concrete details: flower species, fabrics, lighting temperature, composition cues.
- No people, vehicles, or structural edits.
- Output format: line 1 contains MODE label, next a single paragraph prompt.
"""

def strip_env(value: Optional[str]) -> str:
    if not value:
        return ""
    return value.strip().strip('"').strip("'")


def detect_mode(user_prompt: str) -> str:
    """Mode A: element placement, Mode B: full venue decoration."""
    p = (user_prompt or "").lower().strip()

    mode_a_triggers = [
        "add mandap", "place mandap", "add arch", "place arch", "add backdrop",
        "place backdrop", "only florals", "just flowers", "only fabric",
        "just drapery", "aisle decor", "entrance decor", "stage setup", "add a backdrop",
    ]
    if any(trigger in p for trigger in mode_a_triggers):
        return "A"

    mode_b_triggers = [
        "decorate complete", "full decoration", "decorate the entire", "decorate entire",
        "decorate like this", "decorate hall", "haldi", "mehndi", "sangeet",
        "reception", "nikah", "wedding decoration", "transform", "fill the venue",
        "complete setup", "complete decoration", "same decoration", "this style",
        "this theme", "like this", "decorate for", "apply this", "make it look",
        "same vibe", "same feel", "full setup", "lavender", "purple theme",
        "yellow theme", "decorate with", "coral", "minimal decoration",
        "grand decoration", "heavy decoration", "opulent", "entire space",
        "whole hall", "whole venue",
    ]
    if any(trigger in p for trigger in mode_b_triggers):
        return "B"

    if any(word in p for word in ["decorate", "decoration", "setup", "theme"]):
        return "B"

    return "B"


def normalize_base64(image_value: Optional[str]) -> Optional[str]:
    if not image_value:
        return None
    value = image_value.strip()
    if value.startswith("data:") and "," in value:
        return value.split(",", 1)[1]
    return value


async def fetch_image_as_base64(url: str) -> Optional[str]:
    if not url:
        return None
    try:
        async with httpx.AsyncClient(timeout=30.0, follow_redirects=True) as client:
            response = await client.get(url)
            if response.status_code != 200:
                logger.warning(f"Image URL fetch failed ({response.status_code}): {url[:120]}")
                return None
            if not response.content:
                logger.warning(f"Image URL returned empty content: {url[:120]}")
                return None
            return base64.b64encode(response.content).decode("utf-8")
    except Exception as fetch_error:
        logger.warning(f"Image URL fetch exception: {fetch_error}")
        return None


def build_prompt(req: GenerateRequest) -> str:
    base_prompt = (req.prompt or "").strip()
    parts = [f"Client request: {base_prompt}" if base_prompt else "Client request: design the venue decor."]
    
    event_prompts = {
        "Ultra-Luxury Wedding": "Aesthetic style: Opulent ultra-luxury wedding design. Rich, dense floral arrangements using cream roses, pastel hydrangeas, crystal hanging chandeliers, ambient warm candle uplighting, luxury seating with gold accents, grand tablescapes with fine silverware, fairy-tale atmospheric romantic look.",
        "Small Function": "Aesthetic style: Intimate small function design suitable for birthday parties, success celebrations, or private gatherings. Festive balloon arches, warm fairy lights, colorful and elegant banners, cozy lounge seating, modern geometric backdrops, warm and cheerful ambient lighting, welcoming and elegant atmosphere.",
        "Corporate Conference": "Aesthetic style: Professional high-end corporate conference design. Clean geometric stage, neat rows of professional seating, modern LED presentation panels, sleek linear ceiling lighting, minimalist business branding banners, neat white floral accents, sophisticated corporate look.",
        "Global Exhibition": "Aesthetic style: Modern international exhibition hall. Premium display booths, high-tech spotlight tracks, sleek industrial display stands, clean geometric partitions, brand banners, polished minimal design, bright daylight studio lighting.",
        "Fashion Show": "Aesthetic style: High-fashion runway catwalk design. Sleek raised runway platform, dramatic overhead theatrical spotlights, sharp directional key lighting, dark textured backdrop panels, front-row premium seating, high-contrast dramatic shadows.",
        "Product Launch": "Aesthetic style: Premium product launch stage design. Elegant display pedestal at the focal point, dramatic backlighting, clean sophisticated backdrop panel, sleek modern presentation lounge furniture, professional branding presentation layout.",
        "Cultural Festival": "Aesthetic style: Vibrant festive cultural celebration. Ornate traditional drapery, rich ethnic patterns, dense marigold and orchid floral garlands, hanging brass lamps, traditional lanterns, colorful and celebratory warm ambient lighting."
    }

    if req.function_type:
        parts.append(f"Event type: {req.function_type}")
        if req.function_type in event_prompts:
            parts.append(event_prompts[req.function_type])
    if req.theme:
        parts.append(f"Theme: {req.theme}")
    if req.space:
        parts.append(f"Venue space: {req.space} at Fairmont Mumbai")
    parts.extend([
        "Output goal: one photorealistic transformed venue image suitable for client presentation.",
        "Hard constraints: preserve architecture, camera angle, and structural geometry.",
        "Decor policy: additive decor only; do not alter walls, floor, ceiling, pillars, doors, or windows.",
        "Scene policy: no people, no text overlays, no logos, no watermarks.",
    ])
    return "\n".join(parts)


def build_flux_prompt(user_prompt: str, mode: str) -> str:
    mode_instruction = (
        "MODE A: place only the requested decor elements while preserving all venue structure."
        if mode == "A"
        else "MODE B: fully decorate the venue in the reference style while preserving architecture."
    )
    return f"""You are a photorealistic wedding decor compositor.

{mode_instruction}

Input semantics:
- Image 1 is the real venue to preserve.
- Image 2 is the style reference to transfer.

CRITICAL OVERLAY REQUIREMENT:
The venue space is exactly provided in Image 1. You must keep every single structural pixel, wall, ceiling, window, door, light fixture, pillar, floor pattern, and camera perspective 100% identical and unchanged. ONLY add/overlay the decor elements (such as tables, chairs, drapery, flowers, table settings) in the empty spaces or overlaying the tables. Do not alter the background or structure of the venue in any way.

Generation requirements:
- Keep venue architecture, perspective, lens feel, and room proportions intact.
- Add decor with realistic scale, depth, shadows, reflections, and occlusion.
- Use premium wedding materials and cohesive styling (florals, fabric, lighting, tablescape).
- Maintain visual plausibility for walkways and seating clearances.
- Keep color harmony consistent with the requested theme.

Do not:
- Change walls/floor/ceiling/pillars/windows/doors geometry.
- Add people, signage text, logos, or watermarks.
- Produce cartoonish, painterly, or low-detail output.

User requirements:
{user_prompt}
"""


def guess_mime_from_base64(image_base64: str) -> str:
    snippet = (image_base64 or "")[:32]
    if snippet.startswith("iVBORw0KGgo"):
        return "image/png"
    if snippet.startswith("/9j/"):
        return "image/jpeg"
    if snippet.startswith("UklGR"):
        return "image/webp"
    return "image/jpeg"


def base64_to_data_uri(image_base64: str) -> str:
    mime = guess_mime_from_base64(image_base64)
    return f"data:{mime};base64,{image_base64}"


def extract_message_text(content) -> str:
    if isinstance(content, str):
        return content.strip()
    if isinstance(content, list):
        parts = []
        for item in content:
            if isinstance(item, dict) and item.get("type") == "text":
                parts.append(item.get("text", ""))
        return "\n".join(part for part in parts if part).strip()
    return ""


def extract_final_prompt(raw_output: str, mode: str) -> str:
    text = (raw_output or "").strip()
    if not text:
        return ""

    lines = text.split("\n")
    cleaned = []
    for line in lines:
        s = line.strip().upper()
        if s in {
            "MODE A", "MODE B", "MODE: A", "MODE: B", "**MODE A**", "**MODE B**",
            "DETECTED MODE: A", "DETECTED MODE: B", "A", "B"
        }:
            continue
        cleaned.append(line)
    text = "\n".join(cleaned).strip()

    markers = [
        "**FINAL PROMPT:**", "FINAL PROMPT:", "**FINAL PROMPT**", "FINAL PROMPT",
        "**PROMPT:**", "PROMPT:"
    ]
    for marker in markers:
        idx = text.find(marker)
        if idx != -1:
            candidate = text[idx + len(marker):].strip().strip("-").strip()
            if candidate.startswith('"') and candidate.endswith('"'):
                candidate = candidate[1:-1].strip()
            if len(candidate.split()) > 50:
                return candidate

    for starter in [
        "A photorealistic photograph of", "A photorealistic photo of", "A photorealistic image of"
    ]:
        idx = text.find(starter)
        if idx != -1:
            candidate = text[idx:].strip()
            if len(candidate.split()) > 80:
                return candidate

    longest = ""
    for para in text.split("\n\n"):
        p = para.strip()
        if p and len(p.split()) > len(longest.split()) and p[0] not in "#*-=":
            longest = p
    if len(longest.split()) > 80:
        return longest

    min_wc = 280 if mode == "A" else 340
    prose_lines = []
    for line in text.split("\n"):
        s = line.strip()
        if not s:
            continue
        if s[0] in "#*-=":
            continue
        if ":" in s and len(s.split()) < 8:
            continue
        if len(s.split()) > 8:
            prose_lines.append(s)
    combined = " ".join(prose_lines).strip()
    if len(combined.split()) > min_wc:
        return combined

    return text


def build_analysis_instruction(mode: str, user_prompt: str) -> str:
    if mode == "B":
        return f"""Analyze both images for FULL VENUE DECORATION.
IMAGE 1 (VENUE): room type, floor material/pattern, walls/panels/trim, ceiling structure, lighting, perspective, top venue identity features.
IMAGE 2 (DECOR): event style, color palette, ceiling decor, wall decor, floor treatment, focal point, seating, florals (species), density, distribution.
USER REQUEST: \"{user_prompt}\"
Return concise structured analysis."""
    return f"""Analyze both images for ELEMENT PLACEMENT.
IMAGE 1 (VENUE): floor, walls, ceiling, lighting, perspective, top venue identity features.
IMAGE 2 (DECOR): decor element types, approximate dimensions, floral details, fabrics, top colors.
USER REQUEST: \"{user_prompt}\"
Return concise structured analysis."""


async def generate_prompt_with_groq(
    user_prompt: str,
    venue_image_base64: str,
    decor_image_base64: str,
    groq_api_key: str,
    groq_model: str,
) -> tuple[str, str, str]:
    mode = detect_mode(user_prompt)
    venue_data_uri = base64_to_data_uri(venue_image_base64)
    decor_data_uri = base64_to_data_uri(decor_image_base64)

    headers = {
        "Authorization": f"Bearer {groq_api_key}",
        "Content-Type": "application/json",
    }
    groq_url = "https://api.groq.com/openai/v1/chat/completions"

    analysis_instruction = build_analysis_instruction(mode, user_prompt)

    async with httpx.AsyncClient(timeout=90.0) as client:
        analysis_payload = {
            "model": groq_model,
            "messages": [
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": analysis_instruction},
                        {"type": "image_url", "image_url": {"url": venue_data_uri}},
                        {"type": "image_url", "image_url": {"url": decor_data_uri}},
                    ],
                }
            ],
            "temperature": 0.05,
            "max_tokens": 2000,
        }
        analysis_resp = await client.post(groq_url, headers=headers, json=analysis_payload)
        analysis_resp.raise_for_status()
        analysis_json = analysis_resp.json()
        analysis_content = analysis_json.get("choices", [{}])[0].get("message", {}).get("content", "")
        analysis = extract_message_text(analysis_content)

        prompt_payload = {
            "model": groq_model,
            "messages": [
                {"role": "system", "content": SYSTEM_PROMPT_ADVANCED_V2},
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": f"""Analysis:
---
{analysis}
---
User: \"{user_prompt}\"
Output MODE on line 1, then the prompt paragraph. Nothing else.""",
                        },
                        {"type": "image_url", "image_url": {"url": venue_data_uri}},
                        {"type": "image_url", "image_url": {"url": decor_data_uri}},
                    ],
                },
            ],
            "temperature": 0.2,
            "max_tokens": 3000,
        }
        prompt_resp = await client.post(groq_url, headers=headers, json=prompt_payload)
        prompt_resp.raise_for_status()
        prompt_json = prompt_resp.json()
        raw_content = prompt_json.get("choices", [{}])[0].get("message", {}).get("content", "")
        raw_output = extract_message_text(raw_content)

        final_prompt = extract_final_prompt(raw_output, mode)
        min_words = 280 if mode == "A" else 340

        if len(final_prompt.split()) < min_words:
            retry_payload = {
                "model": groq_model,
                "messages": [
                    {"role": "system", "content": SYSTEM_PROMPT_ADVANCED_V2},
                    {
                        "role": "user",
                        "content": [
                            {
                                "type": "text",
                                "text": (
                                    f"Analysis:\n{analysis}\n\n"
                                    f"User: \"{user_prompt}\"\n\n"
                                    f"MODE + prompt only. {'380-480' if mode == 'B' else '320-420'} words. "
                                    "Start with 'A photorealistic photograph of...'"
                                ),
                            },
                            {"type": "image_url", "image_url": {"url": venue_data_uri}},
                            {"type": "image_url", "image_url": {"url": decor_data_uri}},
                        ],
                    },
                ],
                "temperature": 0.2,
                "max_tokens": 2500,
            }
            try:
                retry_resp = await client.post(groq_url, headers=headers, json=retry_payload)
                retry_resp.raise_for_status()
                retry_json = retry_resp.json()
                retry_content = retry_json.get("choices", [{}])[0].get("message", {}).get("content", "")
                retry_output = extract_message_text(retry_content)
                retry_prompt = extract_final_prompt(retry_output, mode)
                if len(retry_prompt.split()) > len(final_prompt.split()):
                    final_prompt = retry_prompt
            except Exception as retry_error:
                logger.warning(f"Groq retry prompt enhancement failed, using initial Groq prompt: {retry_error}")

    if not final_prompt or len(final_prompt.split()) < 50:
        final_prompt = user_prompt

    return final_prompt, analysis, mode


def resolve_result_image_url(data: dict) -> Optional[str]:
    """Find image URL across multiple response shapes."""
    if not isinstance(data, dict):
        return None

    # Some providers return URL keys at top-level.
    for key in ("sample", "url", "image_url", "output"):
        val = data.get(key)
        if isinstance(val, str) and val.startswith("http"):
            return val

    result = data.get("result")
    if isinstance(result, dict):
        for key in ("sample", "url", "image_url", "output"):
            val = result.get(key)
            if isinstance(val, str) and val.startswith("http"):
                return val
    elif isinstance(result, str) and result.startswith("http"):
        return result

    return None


def guess_mime_from_binary(data: bytes) -> str:
    if data.startswith(b"\x89PNG\r\n\x1a\n"):
        return "image/png"
    if data.startswith(b"\xff\xd8\xff"):
        return "image/jpeg"
    if data.startswith(b"RIFF") and b"WEBP" in data[:16]:
        return "image/webp"
    return "image/png"


def build_nano_banana_contents(prompt: str, venue_image: Optional[str], design_image: Optional[str]) -> list:
    parts = [{"text": prompt}]
    if venue_image:
        parts.append({"inlineData": {"mimeType": guess_mime_from_base64(venue_image), "data": venue_image}})
    if design_image:
        parts.append({"inlineData": {"mimeType": guess_mime_from_base64(design_image), "data": design_image}})
    return [{"role": "user", "parts": parts}]


def extract_nano_banana_image_and_text(data: dict) -> tuple[Optional[str], Optional[str], Optional[str]]:
    """Returns (image_base64, mime_type, text_message)."""
    if not isinstance(data, dict):
        return None, None, None

    candidates = data.get("candidates") or []
    text_fragments = []
    for candidate in candidates:
        content = (candidate or {}).get("content") or {}
        for part in content.get("parts") or []:
            if isinstance(part, dict):
                inline = part.get("inlineData") or part.get("inline_data")
                if isinstance(inline, dict) and inline.get("data"):
                    return inline.get("data"), inline.get("mimeType") or inline.get("mime_type") or "image/png", None
                if part.get("text"):
                    text_fragments.append(str(part.get("text")))

    return None, None, "\n".join(text_fragments).strip() or None


async def generate_with_nano_banana(
    prompt: str,
    api_key: str,
    model: str,
    venue_image: Optional[str],
    design_image: Optional[str],
) -> dict:
    payload = {
        "contents": build_nano_banana_contents(prompt, venue_image, design_image),
        "generationConfig": {
            "responseModalities": ["TEXT", "IMAGE"],
        },
    }

    candidate_models = [
        model,
        "nano-banana-pro-preview",
        "gemini-2.5-flash-image",
        "gemini-3-pro-image-preview",
        "gemini-3.1-flash-image-preview",
    ]

    # Preserve order while removing duplicates.
    seen = set()
    models_to_try = []
    for candidate in candidate_models:
        c = (candidate or "").strip()
        if c and c not in seen:
            seen.add(c)
            models_to_try.append(c)

    async with httpx.AsyncClient(timeout=90.0) as client:
        for idx, model_name in enumerate(models_to_try):
            url = f"https://generativelanguage.googleapis.com/v1beta/models/{model_name}:generateContent?key={api_key}"
            response = await client.post(url, json=payload, headers={"Content-Type": "application/json"})
            body_text = response.text.strip()

            if response.status_code == 200:
                try:
                    data = response.json()
                except Exception:
                    return {
                        "success": False,
                        "error": "Nano Banana invalid response",
                        "description": body_text[:1000],
                    }

                image_b64, mime_type, text_message = extract_nano_banana_image_and_text(data)
                if not image_b64:
                    return {
                        "success": False,
                        "error": "Nano Banana returned no image",
                        "description": text_message or body_text[:1000] or "No image found in response",
                    }

                return {
                    "success": True,
                    "image_data": image_b64,
                    "mime_type": mime_type or "image/png",
                    "description": "Design generated successfully",
                    "provider": "nano_banana",
                    "model": model_name,
                }

            has_next_model = idx < len(models_to_try) - 1
            not_found_model = response.status_code == 404 and "not found" in body_text.lower()
            transient_error = response.status_code in {429, 500, 502, 503, 504}
            should_retry = has_next_model and (not_found_model or transient_error)
            if should_retry:
                logger.warning(
                    f"Nano Banana model attempt failed ({model_name}), retrying with next model. "
                    f"Status={response.status_code}, body={body_text[:300]}"
                )
                continue

            logger.error(f"Nano Banana API error {response.status_code}: {body_text[:1000]}")
            return {
                "success": False,
                "error": f"Nano Banana API error: {response.status_code}",
                "description": body_text[:1000] or "Request failed",
            }

    return {
        "success": False,
        "error": "Nano Banana model resolution failed",
        "description": "No compatible Nano Banana model could be used.",
    }


async def generate_flux_variant_from_payload(
    target_url: str,
    api_payload: dict,
    flux_api_key: str,
    max_attempts: int = 100,
) -> dict:
    """Generate a single FLUX image variant from a prepared payload."""
    async with httpx.AsyncClient(timeout=30.0) as client:
        response = await client.post(
            target_url,
            json=api_payload,
            headers={
                "x-key": flux_api_key,
                "Content-Type": "application/json",
            },
        )

    if response.status_code not in [200, 201, 202]:
        body = response.text.strip()
        return {
            "success": False,
            "error": f"Variant submit failed: {response.status_code}",
            "description": body[:500] if body else "Submit failed",
        }

    try:
        data = response.json()
    except Exception:
        return {
            "success": False,
            "error": "Variant submit invalid response",
            "description": response.text[:500],
        }

    task_id = data.get("id", "")
    polling_url = data.get("polling_url", f"https://api.bfl.ai/v1/get_result?id={task_id}")
    result_url = resolve_result_image_url(data)

    if not result_url and not task_id:
        return {
            "success": False,
            "error": "Variant missing task ID or image URL",
            "description": "Flux response did not include polling info",
        }

    if not result_url:
        async with httpx.AsyncClient(timeout=15.0) as client:
            for _ in range(max_attempts):
                await asyncio.sleep(3)
                poll_response = await client.get(
                    polling_url,
                    headers={"x-key": flux_api_key, "Content-Type": "application/json"},
                )
                if poll_response.status_code != 200:
                    continue
                poll_data = poll_response.json()
                status = str(poll_data.get("status", "unknown")).lower()
                if status in ["ready", "succeeded", "complete", "completed"]:
                    result_url = resolve_result_image_url(poll_data)
                    if result_url:
                        break
                    return {
                        "success": False,
                        "error": "Variant result missing image",
                        "description": str(poll_data)[:200],
                    }
                if status in ["error", "failed", "content moderated", "request moderated"]:
                    return {
                        "success": False,
                        "error": "Variant generation failed",
                        "description": str(poll_data)[:200],
                    }

    if not result_url:
        return {
            "success": False,
            "error": "Variant timeout",
            "description": "Variant generation timed out",
        }

    async with httpx.AsyncClient(timeout=60.0) as client:
        img_response = await client.get(result_url)
    if img_response.status_code != 200:
        return {
            "success": False,
            "error": "Variant download failed",
            "description": f"Status: {img_response.status_code}",
        }

    img_base64 = base64.b64encode(img_response.content).decode("utf-8")
    return {
        "success": True,
        "image_data": img_base64,
        "mime_type": "image/png",
        "task_id": task_id,
    }

@api_router.get("/")
async def root():
    return {"message": "Fairmont x LoversAI API"}

@api_router.post("/generate")
async def generate_image(req: GenerateRequest):
    try:
        user_prompt = build_prompt(req)
        mode = detect_mode(user_prompt)
        logger.info(f"Generating with prompt: {user_prompt[:100]}...")
        flux_prompt = build_flux_prompt(user_prompt, mode)
        
        # Handle image inputs - support base64 + URL fallback + legacy reference_image
        venue_image = normalize_base64(req.venue_image)
        design_image = normalize_base64(req.design_image)

        if not venue_image and req.venue_image_url:
            logger.info("Venue base64 missing, fetching from venue_image_url...")
            venue_image = await fetch_image_as_base64(req.venue_image_url)

        if not design_image and req.design_image_url:
            logger.info("Design base64 missing, fetching from design_image_url...")
            design_image = await fetch_image_as_base64(req.design_image_url)
        
        # Legacy support: if reference_image exists, treat it as design_image
        if req.reference_image and not design_image:
            design_image = normalize_base64(req.reference_image)
            logger.info("Using legacy reference_image as design_image")
        
        has_venue = bool(venue_image)
        has_design = bool(design_image)
        requested_variants = req.variant_count if req.variant_count is not None else (3 if req.high_quality else 1)
        variant_count = max(1, min(3, int(requested_variants or 1)))
        requested_input_mode = "dual_image" if (has_venue and has_design) else ("venue_only" if has_venue else "design_only" if has_design else "text_only")
        effective_input_mode = requested_input_mode
        prompt_source = "rule_based"
        groq_analysis_excerpt = ""
        
        logger.info(f"=== API RECEIVED ===")
        logger.info(f"prompt: {user_prompt[:50]}...")
        logger.info(f"function_type: {req.function_type}")
        logger.info(f"space: {req.space}")
        logger.info(f"venue_image provided: {bool(venue_image)}, length: {len(venue_image) if venue_image else 0}")
        logger.info(f"design_image provided: {bool(design_image)}, length: {len(design_image) if design_image else 0}")
        logger.info(f"reference_image provided: {bool(req.reference_image)}, length: {len(req.reference_image) if req.reference_image else 0}")
        logger.info(f"venue_image_url provided: {bool(req.venue_image_url)}")
        logger.info(f"design_image_url provided: {bool(req.design_image_url)}")
        logger.info(f"variant_count requested: {variant_count}")
        
        # =========================================================
        # Provider selection
        # Prefer Nano Banana (Gemini image generation) when configured.
        # =========================================================

        nano_banana_api_key = (
            strip_env(os.environ.get("NANO_BANANA_API_KEY", ""))
            or strip_env(os.environ.get("GEMINI_API_KEY", ""))
        )
        nano_banana_model = os.environ.get("NANO_BANANA_MODEL", "nano-banana-pro-preview")

        # Backward-compatible convenience: if FLUX key is a Gemini-style key, use it for Nano Banana.
        flux_key_for_detection = strip_env(os.environ.get("FLUX_API_KEY", ""))
        if not nano_banana_api_key and flux_key_for_detection.startswith("AIza"):
            nano_banana_api_key = flux_key_for_detection

        use_nano_banana = bool(nano_banana_api_key) and (has_venue or has_design)
        if use_nano_banana:
            nano_tasks = [
                generate_with_nano_banana(
                    prompt=flux_prompt,
                    api_key=nano_banana_api_key,
                    model=nano_banana_model,
                    venue_image=venue_image,
                    design_image=design_image,
                )
                for _ in range(variant_count)
            ]
            nano_results = await asyncio.gather(*nano_tasks)
            nano_variants = []
            nano_failure = None
            fallback_to_flux = False

            for nano_result in nano_results:
                if nano_result.get("success"):
                    nano_variants.append({
                        "image_data": nano_result.get("image_data"),
                        "mime_type": nano_result.get("mime_type", "image/png"),
                        "provider": "nano_banana",
                        "model": nano_result.get("model", nano_banana_model),
                    })
                    continue

                nano_failure = nano_result
                nano_error = str(nano_result.get("error", ""))
                nano_description = str(nano_result.get("description", ""))
                fallback_to_flux = fallback_to_flux or (
                    nano_error == "Nano Banana returned no image"
                    and ("provide the images" in nano_description.lower() or "provide images" in nano_description.lower())
                )

            if nano_variants:
                primary = nano_variants[0]
                return {
                    "success": True,
                    "image_data": primary.get("image_data"),
                    "mime_type": primary.get("mime_type", "image/png"),
                    "description": "Design generated successfully",
                    "provider": "nano_banana",
                    "model": primary.get("model", nano_banana_model),
                    "input_mode": requested_input_mode,
                    "effective_input_mode": requested_input_mode,
                    "mode": "element" if mode == "A" else "full_decoration",
                    "prompt_source": "nano_banana",
                    "prompt_words": len(flux_prompt.split()),
                    "variants": nano_variants,
                    "variants_requested": variant_count,
                    "variants_generated": len(nano_variants),
                }
            if not fallback_to_flux:
                logger.warning(f"Nano Banana path failed: {str((nano_failure or {}).get('description', ''))[:200]}")
                return nano_failure or {"success": False, "error": "Nano Banana failed", "description": "No successful variant"}

        # =========================================================
        # FLUX API INTEGRATION - Black Forest Labs
        # Model: flux-kontext-pro (image editing model)
        # =========================================================
        
        flux_api_key = strip_env(os.environ.get("FLUX_API_KEY", ""))
        flux_api_url = os.environ.get("FLUX_API_URL", "https://api.bfl.ai/v1/flux-kontext-pro")
        flux_text_api_url = os.environ.get("FLUX_TEXT_API_URL", "https://api.bfl.ai/v1/flux-2-max")
        groq_api_key = strip_env(os.environ.get("GROQ_API_KEY", "")) or strip_env(os.environ.get("EMERGENT_LLM_KEY", ""))
        groq_vision_model = os.environ.get("GROQ_VISION_MODEL", "meta-llama/llama-4-scout-17b-16e-instruct")
        if groq_api_key and not groq_api_key.startswith("gsk_"):
            logger.warning("GROQ key format appears invalid; skipping Groq pipeline and using Flux fallback path.")
            groq_api_key = ""
        
        if not flux_api_key:
            return {
                "success": False,
                "error": "FLUX_API_KEY not configured. Please set FLUX_API_KEY in .env file.",
                "description": "API key missing"
            }

        # High-priority path from user requirements:
        # Use both venue+decor images to build a descriptive prompt with Groq Vision,
        # then generate via flux-2-max text-to-image.
        if has_venue and has_design and groq_api_key:
            try:
                logger.info("Running Groq dual-image analysis for prompt generation...")
                groq_prompt, analysis, detected_mode = await generate_prompt_with_groq(
                    user_prompt=user_prompt,
                    venue_image_base64=venue_image,
                    decor_image_base64=design_image,
                    groq_api_key=groq_api_key,
                    groq_model=groq_vision_model,
                )
                flux_prompt = groq_prompt
                mode = detected_mode
                prompt_source = "groq_vision_dual_image"
                effective_input_mode = "groq_text_from_dual_image"
                groq_analysis_excerpt = analysis[:600]
                logger.info(f"Groq prompt generated: {len(flux_prompt.split())} words")
            except Exception as groq_error:
                logger.warning(f"Groq pipeline failed, using Flux fallback path: {groq_error}")
        elif has_venue and has_design and not groq_api_key:
            logger.warning("GROQ_API_KEY missing. Falling back to Flux image endpoint without Groq prompt analysis.")

        # If Groq dual-image prompt is available, prefer text endpoint.
        if prompt_source == "groq_vision_dual_image":
            target_url = flux_text_api_url
        else:
            target_url = flux_api_url if (has_venue or has_design) else flux_text_api_url
        is_text_endpoint = "flux-2-max" in target_url
        
        # Prepare request payload.
        if is_text_endpoint:
            api_payload = {
                "prompt": flux_prompt,
                "width": 1024,
                "height": 768,
                "output_format": "png",
                "safety_tolerance": 2,
            }
        else:
            api_payload = {
                "prompt": flux_prompt,
                "width": 1024,
                "height": 768,
            }
        
        # Add images based on availability
        if not is_text_endpoint and has_venue and has_design:
            # Dual image input - venue + design inspiration
            api_payload["input_image"] = venue_image
            api_payload["image_prompt"] = design_image
            logger.info("Using dual image-to-image mode (venue + design)")
        elif not is_text_endpoint and has_venue:
            # Single venue image only
            api_payload["input_image"] = venue_image
            logger.info("Using single image mode (venue only)")
        elif not is_text_endpoint and has_design:
            # Single design image only
            api_payload["input_image"] = design_image
            logger.info("Using single image mode (design only)")
        elif is_text_endpoint:
            logger.info("Using text-to-image mode (no input images)")
        
        # Submit request to FLUX API
        async with httpx.AsyncClient(timeout=30.0) as client:
            response = await client.post(
                target_url,
                json=api_payload,
                headers={
                    "x-key": flux_api_key,  # Note: lowercase 'x-key'
                    "Content-Type": "application/json"
                }
            )
            
            logger.info(f"FLUX API status: {response.status_code}")
            response_text = response.text.strip()
            if response_text:
                logger.info(f"FLUX API response body: {response_text[:1000]}")
            
            # Handle error responses
            if response.status_code == 401:
                return {"success": False, "error": "Invalid API key", "description": response_text or "Check FLUX_API_KEY in .env"}
            elif response.status_code == 402:
                return {"success": False, "error": "Insufficient credits", "description": response_text or "Top up your BFL account"}
            elif response.status_code == 422:
                # Fallback 1: retry without image_prompt when dual-image input fails validation.
                if not is_text_endpoint and "image_prompt" in api_payload:
                    logger.warning("422 error - retrying without image_prompt...")
                    fallback_payload = dict(api_payload)
                    fallback_payload.pop("image_prompt", None)
                    response = await client.post(
                        target_url,
                        json=fallback_payload,
                        headers={
                            "x-key": flux_api_key,
                            "Content-Type": "application/json"
                        }
                    )
                    response_text = response.text.strip()
                    if response_text:
                        logger.info(f"FLUX fallback response body: {response_text[:1000]}")
                    if response.status_code in [200, 201, 202]:
                        api_payload = fallback_payload
                    logger.info(f"Fallback status: {response.status_code}")

                # Fallback 2: if image endpoint still fails, fall back to text endpoint.
                if response.status_code == 422 and (has_venue or has_design):
                    logger.warning("422 persists - switching to flux-2-max text-to-image fallback...")
                    text_payload = {
                        "prompt": flux_prompt,
                        "width": 1440,
                        "height": 960,
                        "output_format": "png",
                        "safety_tolerance": 2,
                    }
                    response = await client.post(
                        flux_text_api_url,
                        json=text_payload,
                        headers={
                            "x-key": flux_api_key,
                            "Content-Type": "application/json"
                        }
                    )
                    response_text = response.text.strip()
                    if response_text:
                        logger.info(f"FLUX text fallback response body: {response_text[:1000]}")
                    if response.status_code in [200, 201, 202]:
                        target_url = flux_text_api_url
                        is_text_endpoint = True
                        effective_input_mode = "text_fallback"
                        api_payload = text_payload
                    logger.info(f"Text fallback status: {response.status_code}")
            elif response.status_code == 429:
                return {"success": False, "error": "Rate limited", "description": response_text or "Wait and try again"}
            elif response.status_code not in [200, 201, 202]:
                return {"success": False, "error": f"API error: {response.status_code}", "description": response_text[:1000] or response.text[:200]}
            
            # Get task ID and polling URL from response
            data = response.json()
            task_id = data.get("id", "")
            polling_url = data.get("polling_url", f"https://api.bfl.ai/v1/get_result?id={task_id}")
            result_url = resolve_result_image_url(data)
            
            logger.info(f"Task ID: {task_id}, Polling for result...")

            extra_variant_tasks = []
            if variant_count > 1:
                extra_variant_tasks = [
                    asyncio.create_task(
                        generate_flux_variant_from_payload(
                            target_url=target_url,
                            api_payload=api_payload,
                            flux_api_key=flux_api_key,
                        )
                    )
                    for _ in range(variant_count - 1)
                ]
        
        # =========================================================
        # POLL FOR RESULT (async polling)
        # =========================================================
        
        max_attempts = 100  # Max 5 minutes (100 * 3s)
        # Some endpoints may return direct URL without polling.
        if not result_url and not task_id:
            return {
                "success": False,
                "error": "No task ID or image URL",
                "description": "Flux response did not include polling info"
            }
        
        if not result_url:
            async with httpx.AsyncClient(timeout=15.0) as client:
                for attempt in range(1, max_attempts + 1):
                    await asyncio.sleep(3)  # Poll every 3 seconds
                    
                    poll_response = await client.get(
                        polling_url,
                        headers={"x-key": flux_api_key, "Content-Type": "application/json"}
                    )
                    
                    if poll_response.status_code != 200:
                        logger.warning(f"Poll attempt {attempt}: status {poll_response.status_code}")
                        continue
                    
                    poll_data = poll_response.json()
                    status = str(poll_data.get("status", "unknown")).lower()
                    
                    logger.info(f"Poll {attempt}: status = {status}")
                    
                    if status in ["ready", "succeeded", "complete", "completed"]:
                        result_url = resolve_result_image_url(poll_data)
                        if result_url:
                            logger.info(f"Result ready after {attempt * 3}s")
                            break
                        return {"success": False, "error": "No image in result", "description": str(poll_data)[:200]}
                    
                    if status in ["error", "failed", "content moderated", "request moderated"]:
                        return {"success": False, "error": "Generation failed", "description": str(poll_data)[:200]}
                
                if not result_url:
                    return {"success": False, "error": "Timeout", "description": "Generation timed out after 5 minutes"}
        
        # =========================================================
        # DOWNLOAD OUTPUT IMAGE
        # =========================================================
        
        async with httpx.AsyncClient(timeout=60.0) as client:
            img_response = await client.get(result_url)
            
            if img_response.status_code != 200:
                return {"success": False, "error": "Failed to download image", "description": f"Status: {img_response.status_code}"}
            
            # Convert to base64
            img_base64 = base64.b64encode(img_response.content).decode("utf-8")
            variants = [{"image_data": img_base64, "mime_type": "image/png", "task_id": task_id}]

            if variant_count > 1 and extra_variant_tasks:
                done_tasks, pending_tasks = await asyncio.wait(extra_variant_tasks, timeout=5.0)
                for pending_task in pending_tasks:
                    pending_task.cancel()
                for completed_task in done_tasks:
                    try:
                        extra = completed_task.result()
                    except Exception as extra_error:
                        logger.warning(f"Variant generation task raised: {extra_error}")
                        continue
                    if extra.get("success"):
                        variants.append({
                            "image_data": extra.get("image_data"),
                            "mime_type": extra.get("mime_type", "image/png"),
                            "task_id": extra.get("task_id", ""),
                        })
                    else:
                        logger.warning(f"Variant generation attempt failed: {extra.get('description', '')[:200]}")
            
            return {
                "success": True,
                "image_data": variants[0]["image_data"],
                "mime_type": "image/png",
                "description": "Design generated successfully",
                "input_mode": requested_input_mode,
                "effective_input_mode": effective_input_mode,
                "task_id": task_id,
                "mode": "element" if mode == "A" else "full_decoration",
                "endpoint_used": target_url,
                "prompt_source": prompt_source,
                "prompt_words": len(flux_prompt.split()),
                "groq_analysis_excerpt": groq_analysis_excerpt,
                "variants": variants,
                "variants_requested": variant_count,
                "variants_generated": len(variants),
            }
        
    except httpx.TimeoutException:
        logger.error("FLUX API timeout")
        return {"success": False, "error": "Request timed out", "description": "Try again later"}
    except Exception as e:
        logger.error(f"Generation error: {str(e)}")
        return {"success": False, "error": str(e)}

@api_router.post("/moodboard/save")
async def save_to_moodboard(req: MoodboardSaveRequest):
    if db is None:
        return {"success": False, "error": "Database not available", "id": None}
    doc = {
        "id": str(uuid.uuid4()),
        "session_id": req.session_id,
        "image_data": req.image_data,
        "prompt": req.prompt,
        "filters": req.filters,
        "created_at": datetime.now(timezone.utc).isoformat()
    }
    await db.moodboard_images.insert_one(doc)
    return {"success": True, "id": doc["id"]}

@api_router.get("/moodboard/{session_id}")
async def get_moodboard(session_id: str):
    if db is None:
        return {"images": [], "error": "Database not available"}
    images = await db.moodboard_images.find(
        {"session_id": session_id}, {"_id": 0}
    ).to_list(100)
    return {"images": images}

@api_router.post("/moodboard/download-pdf")
async def download_pdf(req: DownloadRequest):
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import inch
    from reportlab.pdfgen import canvas as pdf_canvas
    from reportlab.lib.utils import ImageReader

    buffer = io.BytesIO()
    c = pdf_canvas.Canvas(buffer, pagesize=A4)
    width, height = A4

    # Title page
    c.setFont("Helvetica-Bold", 28)
    c.drawCentredString(width / 2, height - 2 * inch, "Fairmont x LoversAI")
    c.setFont("Helvetica", 16)
    c.drawCentredString(width / 2, height - 2.6 * inch, "Moodboard Collection")
    c.setFont("Helvetica-Oblique", 11)
    c.drawCentredString(width / 2, height - 3.2 * inch, datetime.now(timezone.utc).strftime("%B %d, %Y"))
    c.showPage()

    for i, img_item in enumerate(req.images):
        try:
            img_bytes = base64.b64decode(img_item["image_data"])
            img_reader = ImageReader(io.BytesIO(img_bytes))

            img_w = width - 2 * inch
            img_h = img_w * 0.65
            x = inch
            y = height - img_h - 1.5 * inch

            c.drawImage(img_reader, x, y, width=img_w, height=img_h, preserveAspectRatio=True)

            c.setFont("Helvetica", 10)
            prompt_text = img_item.get("prompt", "")[:80]
            c.drawString(inch, y - 0.3 * inch, f"Design {i + 1}: {prompt_text}")
            c.showPage()
        except Exception as e:
            logger.error(f"PDF image error: {e}")
            continue

    c.save()
    buffer.seek(0)
    return StreamingResponse(
        buffer,
        media_type="application/pdf",
        headers={"Content-Disposition": "attachment; filename=moodboard.pdf"}
    )

@api_router.post("/moodboard/download-ppt")
async def download_ppt(req: DownloadRequest):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    txBox = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Fairmont x LoversAI"
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "Moodboard Collection"
    p2.font.size = Pt(24)
    p2.alignment = PP_ALIGN.CENTER

    for i, img_item in enumerate(req.images):
        try:
            img_bytes = base64.b64decode(img_item["image_data"])
            img_stream = io.BytesIO(img_bytes)

            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(img_stream, Inches(1), Inches(0.5), Inches(11.333), Inches(6))

            txBox = slide.shapes.add_textbox(Inches(1), Inches(6.7), Inches(11), Inches(0.6))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = f"Design {i + 1}: {img_item.get('prompt', '')[:60]}"
            p.font.size = Pt(14)
            p.alignment = PP_ALIGN.LEFT
        except Exception as e:
            logger.error(f"PPT image error: {e}")
            continue

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return StreamingResponse(
        buffer,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": "attachment; filename=moodboard.pptx"}
    )

import json

# Templates endpoints
TEMPLATES_DIR = ROOT_DIR / "data" / "templates"
TEMPLATES_CONFIG = ROOT_DIR / "data" / "templates.json"

@api_router.get("/templates")
async def list_templates():
    try:
        if TEMPLATES_CONFIG.exists():
            with open(TEMPLATES_CONFIG, "r") as f:
                templates = json.load(f)
            # Check which files actually exist
            for t in templates:
                t["available"] = (TEMPLATES_DIR / t["filename"]).exists()
            return {"templates": templates}
        return {"templates": []}
    except Exception as e:
        logger.error(f"Templates error: {e}")
        return {"templates": []}

@api_router.get("/templates/download/{filename}")
async def download_template(filename: str):
    filepath = TEMPLATES_DIR / filename
    if not filepath.exists():
        return {"error": "Template file not found"}
    return StreamingResponse(
        open(filepath, "rb"),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename={filename}"}
    )

app.include_router(api_router)

app.add_middleware(
    CORSMiddleware,
    allow_credentials=True,
    allow_origins=get_cors_origins(),
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.on_event("shutdown")
async def shutdown_db_client():
    if client is not None:
        client.close()
